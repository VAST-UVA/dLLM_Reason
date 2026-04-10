"""Generate V1.0 overview presentation.

Usage: python docs/generate_ppt.py
Output: docs/dLLM_Reason_V1.0.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID = RGBColor(0x22, 0x22, 0x3A)
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT2 = RGBColor(0x56, 0x9C, 0xD6)
ACCENT3 = RGBColor(0xDC, 0xDC, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xE0, 0x8A, 0x40)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=16, color=LIGHT_GRAY, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.space_before = Pt(space_before)
    return p


def add_bullet(tf, text, font_size=14, color=LIGHT_GRAY, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.level = level
    p.space_before = Pt(3)
    return p


def add_rect(slide, left, top, width, height, fill_color, text="", font_size=11, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Segoe UI"
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 1.5, 1.5, 10, 1.2, "dLLM-Reason", 48, ACCENT, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.8, 10, 0.8,
             "DAG-Guided Discrete Diffusion Language Models for Reasoning",
             24, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10, 0.5, "V1.0 Release", 20, ACCENT3, False, PP_ALIGN.CENTER)

tf = add_text_box(slide, 2.5, 5.2, 8, 1.2, "", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
tf.paragraphs[0].text = "69 Python files  |  ~9,200 lines of code  |  31 config files"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].font.name = "Segoe UI"
add_para(tf, "4 dLLM models  |  5 schedulers  |  4 search methods  |  DAG library system",
         14, LIGHT_GRAY, space_before=4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: Core Insight
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "Core Insight", 32, ACCENT, True)

tf = add_text_box(slide, 0.8, 1.3, 11.5, 1.0, "", 18)
tf.paragraphs[0].text = "Discrete diffusion LMs (dLLMs) generate text by iteratively unmasking tokens."
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.name = "Segoe UI"
add_para(tf, "Standard: random or confidence-based unmasking order.", 18, LIGHT_GRAY)
add_para(tf, "Our approach: use DAG topology to enforce reasoning dependencies.", 18, ACCENT3, True)
add_para(tf, "\"Position j can only unmask after all its DAG parents are unmasked.\"", 16, ORANGE, space_before=12)

# Architecture boxes
y_base = 3.8
add_rect(slide, 1.0, y_base, 3.2, 1.0, RGBColor(0x2D, 0x4A, 0x6E),
         "Model Layer\nMDLM | SEDD | D3PM | LLaDA\n(what to predict)", 12)
add_rect(slide, 5.0, y_base, 3.2, 1.0, RGBColor(0x4A, 0x2D, 0x6E),
         "Scheduler Layer\nRandom | Confidence | DAG\n(where to unmask)", 12)
add_rect(slide, 9.0, y_base, 3.2, 1.0, RGBColor(0x2D, 0x6E, 0x4A),
         "DAG Layer\nTokenDAG + Templates\n(dependency structure)", 12)

# Arrows
for x in [4.25, 8.25]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y_base + 0.3),
                                   Inches(0.7), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

tf2 = add_text_box(slide, 0.8, 5.3, 11.5, 1.5, "", 14)
tf2.paragraphs[0].text = "Key Design: DAG constraints inject at the scheduler layer."
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.color.rgb = ACCENT
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.name = "Segoe UI"
add_para(tf2, "Models need zero modification. Any dLLM + any scheduler = free composition.", 14, LIGHT_GRAY)
add_para(tf2, "ready = (~adj | is_unmasked).all(dim=1)   # one line of GPU tensor math", 13,
         ACCENT3, space_before=8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Version Timeline
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "Version Timeline", 32, ACCENT, True)

versions = [
    ("V0.1", "Foundation", "DiffusionLM base, MDLM, RandomScheduler,\nSampler, Trainer, GSM8K loader", RGBColor(0x3A, 0x5A, 0x7A)),
    ("V0.2", "Baselines", "SEDD, D3PM, Confidence/Linear schedulers,\nMetrics, MultiSchedulerComparison", RGBColor(0x3A, 0x6A, 0x5A)),
    ("V0.3", "DAG Core", "TokenDAG, 6 templates, DAGScheduler,\nDAGSampler, acyclicity constraints", RGBColor(0x5A, 0x3A, 0x7A)),
    ("V0.4", "Search", "Evolutionary, Greedy, RL Policy, NOTEARS,\nFitness functions, DAG analysis", RGBColor(0x7A, 0x5A, 0x3A)),
    ("V0.5", "Training+", "LLaDA-8B, DAG-aware training,\nDiffu-GRPO, 4 benchmark evaluators", RGBColor(0x3A, 0x5A, 0x5A)),
    ("V0.6", "Configs", "31 YAMLs, 2 notebooks,\n3 test suites, scripts", RGBColor(0x5A, 0x5A, 0x3A)),
    ("V0.7", "Library", "SQLite+FAISS store, 3 retrieval channels,\n4 fusion, 3 feedback, merge, fitness", RGBColor(0x6A, 0x3A, 0x5A)),
    ("V1.0", "Integration", "Search↔Library pipeline, CLI,\n30+ library tests, bug fixes", RGBColor(0x2D, 0x6E, 0x4A)),
]

for i, (ver, title, desc, color) in enumerate(versions):
    col = i % 4
    row = i // 4
    x = 0.6 + col * 3.15
    y = 1.3 + row * 3.0

    add_rect(slide, x, y, 2.9, 0.5, color, f"{ver}  {title}", 13, WHITE)
    add_text_box(slide, x + 0.1, y + 0.55, 2.7, 2.0, desc, 10, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: Models & Schedulers
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "Models & Schedulers", 32, ACCENT, True)

# Models section
add_text_box(slide, 0.8, 1.3, 5, 0.5, "4 Discrete Diffusion Models", 20, ACCENT2, True)
models = [
    ("MDLM", "Absorbing-state continuous-time\nGeometric/linear/cosine noise schedules\nContinuous-time ELBO loss"),
    ("SEDD", "Score-entropy discrete diffusion\nScore-based parameterization\nTransition rate weighting"),
    ("D3PM", "Discrete-time structured transitions\nAbsorbing + uniform types\nHybrid lambda loss"),
    ("LLaDA", "Wraps HuggingFace LLaDA-8B-Instruct\nInference-only with any scheduler\nPrompt encoding + generation"),
]
for i, (name, desc) in enumerate(models):
    y = 2.0 + i * 1.25
    add_rect(slide, 0.8, y, 1.5, 0.4, RGBColor(0x2D, 0x4A, 0x6E), name, 13)
    add_text_box(slide, 2.5, y - 0.05, 4, 1.1, desc, 11, LIGHT_GRAY)

# Schedulers section
add_text_box(slide, 7.2, 1.3, 5, 0.5, "5 Unmasking Schedulers", 20, ACCENT2, True)
scheds = [
    ("Random", "Uniform random from masked positions"),
    ("Confidence", "Highest-confidence positions first"),
    ("Linear", "Left-to-right sequential"),
    ("DAGScheduler", "DAG-constrained eligible positions\n+ sub-strategy: all_ready / confidence_topk / proportional"),
    ("Adaptive", "DAG + confidence-aware\nwith bypass for stuck situations"),
]
for i, (name, desc) in enumerate(scheds):
    y = 2.0 + i * 1.05
    add_rect(slide, 7.2, y, 2.0, 0.4, RGBColor(0x4A, 0x2D, 0x6E), name, 12)
    add_text_box(slide, 9.4, y - 0.05, 3.5, 0.9, desc, 11, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: TokenDAG & Templates
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "TokenDAG & Templates", 32, ACCENT, True)

tf = add_text_box(slide, 0.8, 1.3, 6, 0.5, "TokenDAG — Core Data Structure", 20, ACCENT2, True)

tf = add_text_box(slide, 0.8, 1.9, 6, 4.0, "", 13)
items = [
    "Boolean adjacency matrix (seq_len x seq_len) on GPU",
    "Edge (i, j) = \"position i must unmask before j\"",
    "ready_positions(): single batched matrix operation",
    "Constructors: no_edges(), linear_chain(), from_edges(), from_levels()",
    "Mutations: add_edges(), remove_edges(), mutate()",
    "Analysis: topological_levels(), depth(), transitive_closure()",
    "Validation: is_valid() checks acyclicity",
    "to_mask_schedule(num_steps): distributes levels across steps",
]
for item in items:
    add_bullet(tf, "  " + item, 13, LIGHT_GRAY)

add_text_box(slide, 7.2, 1.3, 5.5, 0.5, "6 Predefined Templates", 20, ACCENT2, True)
templates = [
    ("Chain-of-Thought", "Sequential reasoning steps\nParallel within each step"),
    ("Answer-First", "Answer tokens unmask first\nReasoning fills in after"),
    ("Skeleton-Detail", "Structure tokens first\nContent tokens second"),
    ("Bidirectional", "Forward + backward passes\nMeet in middle"),
    ("Interleaved", "Alternating position groups"),
    ("Random", "Random edges with given density\nBaseline for comparison"),
]
for i, (name, desc) in enumerate(templates):
    y = 1.9 + i * 0.88
    add_rect(slide, 7.2, y, 2.2, 0.35, RGBColor(0x2D, 0x6E, 0x4A), name, 11)
    add_text_box(slide, 9.6, y - 0.05, 3.2, 0.8, desc, 10, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: DAG Search
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "DAG Structure Search", 32, ACCENT, True)

tf = add_text_box(slide, 0.8, 1.2, 11, 0.8, "", 16)
tf.paragraphs[0].text = "Automatically discover optimal DAG structures that maximize reasoning accuracy."
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.name = "Segoe UI"

methods = [
    ("Evolutionary\nSearch", "Population-based optimization\nTournament selection + crossover\nTopological mutation\nLibrary seed + writeback",
     RGBColor(0x7A, 0x5A, 0x3A)),
    ("Greedy Edge\nSearch", "Add/remove single edges\nKeep best improvement\nPatience-based early stop\nLibrary writeback",
     RGBColor(0x3A, 0x6A, 0x5A)),
    ("RL Policy\nSearch", "DAGPolicyNetwork\nREINFORCE training\nLearns edge probabilities\nEnd-to-end optimization",
     RGBColor(0x5A, 0x3A, 0x7A)),
    ("Differentiable\n(NOTEARS)", "Continuous relaxation\nGradient-based optimization\nAcyclicity constraint via\ntrace exponential",
     RGBColor(0x3A, 0x5A, 0x7A)),
]

for i, (name, desc, color) in enumerate(methods):
    x = 0.6 + i * 3.15
    add_rect(slide, x, 2.2, 2.9, 0.8, color, name, 14)
    add_text_box(slide, x + 0.15, 3.15, 2.6, 2.0, desc, 12, LIGHT_GRAY)

tf2 = add_text_box(slide, 0.8, 5.3, 11, 1.5, "", 14)
tf2.paragraphs[0].text = "Fitness Functions:"
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.color.rgb = ACCENT2
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.name = "Segoe UI"
add_para(tf2, "accuracy_fitness()  — DAGScheduler -> Sampler -> answer extraction -> EM", 13, LIGHT_GRAY)
add_para(tf2, "perplexity_fitness()  |  combined_fitness()  — multi-signal composite", 13, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: DAG Library
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "DAG Library System", 32, ACCENT, True)

tf = add_text_box(slide, 0.8, 1.1, 11, 0.5, "Persistent storage, retrieval, and feedback for DAG structures", 16, WHITE)

# Left column: components
components = [
    ("DAGStore", "SQLite + FAISS\nCRUD, filtered queries, vector search", RGBColor(0x2D, 0x4A, 0x6E)),
    ("Retrieval\n3 Channels", "Semantic (embedding similarity)\nStructural (edit distance / spectral)\nPerformance (benchmark scores)", RGBColor(0x4A, 0x2D, 0x6E)),
    ("Fusion\n4 Strategies", "Weighted | RRF | Max | Voting\nCombines multi-channel rankings", RGBColor(0x2D, 0x6E, 0x4A)),
    ("Feedback\n3 Sources", "Auto (benchmark accuracy)\nHuman (quality ratings)\nElo (tournament ranking)", RGBColor(0x6A, 0x3A, 0x3A)),
    ("Merge\n3 Strategies", "Union | Intersection | Weighted\nAcyclicity enforcement built-in", RGBColor(0x5A, 0x5A, 0x3A)),
]

for i, (name, desc, color) in enumerate(components):
    y = 1.7 + i * 1.1
    add_rect(slide, 0.6, y, 1.8, 0.6, color, name, 11)
    add_text_box(slide, 2.6, y, 4.0, 1.0, desc, 11, LIGHT_GRAY)

# Right column: ablation
add_text_box(slide, 7.5, 1.6, 5, 0.5, "Ablation Configs (7 YAML)", 18, ORANGE, True)
ablations = [
    "default.yaml — all components enabled",
    "disabled.yaml — library completely off (baseline)",
    "no_retrieval.yaml — storage + feedback only",
    "semantic_only.yaml — single retrieval channel",
    "no_elo.yaml — auto feedback only, no Elo",
    "soft_constraint.yaml — soft vs hard DAG enforcement",
    "rrf_fusion.yaml — RRF vs weighted fusion",
]
tf3 = add_text_box(slide, 7.5, 2.2, 5.3, 4.5, "", 12)
for ab in ablations:
    add_bullet(tf3, "  " + ab, 12, LIGHT_GRAY)

add_para(tf3, "", 6, LIGHT_GRAY, space_before=12)
add_para(tf3, "Every component independently toggleable.", 14, ACCENT, True, space_before=8)
add_para(tf3, "Supports isolated variable experiments.", 14, ACCENT3, space_before=4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: Search ↔ Library Integration
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "Search ↔ Library Integration", 32, ACCENT, True)

# Flow diagram with boxes and arrows
boxes = [
    (0.8, 2.0, 2.5, 1.2, "DAG Library\n\nSQLite + FAISS\nHistorical DAGs\n+ Metadata", RGBColor(0x2D, 0x4A, 0x6E)),
    (4.5, 2.0, 3.0, 1.2, "DAG Search\n\nEvolutionary / Greedy\nSeed from library\nEvolve population", RGBColor(0x5A, 0x3A, 0x7A)),
    (8.8, 2.0, 3.0, 1.2, "Evaluation\n\nModel + DAGScheduler\nBenchmark accuracy\nFitness score", RGBColor(0x2D, 0x6E, 0x4A)),
]
for x, y, w, h, text, color in boxes:
    add_rect(slide, x, y, w, h, color, text, 12)

# Arrows
for (x1, x2) in [(3.35, 4.45), (7.55, 8.75)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x1), Inches(2.45),
                                   Inches(x2 - x1), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

# Return arrow (bottom)
tf = add_text_box(slide, 3.5, 3.5, 6, 0.4, "← writeback best DAG + fitness to library ←", 14, ORANGE, True, PP_ALIGN.CENTER)

# Details
tf2 = add_text_box(slide, 0.8, 4.3, 11, 2.8, "", 14)
tf2.paragraphs[0].text = "Seed Phase (Library → Search)"
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.color.rgb = ACCENT2
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.name = "Segoe UI"
add_bullet(tf2, "  Retrieve semantically similar DAGs by task description", 13, LIGHT_GRAY)
add_bullet(tf2, "  Pull top-performing DAGs from performance channel", 13, LIGHT_GRAY)
add_bullet(tf2, "  Fill at most half the initial population from library, rest random", 13, LIGHT_GRAY)

add_para(tf2, "Writeback Phase (Search → Library)", 16, ACCENT2, True, space_before=12)
add_bullet(tf2, "  Best DAG automatically stored with source='search', method, fitness score", 13, LIGHT_GRAY)
add_bullet(tf2, "  Available for future searches — library evolves over time", 13, LIGHT_GRAY)
add_bullet(tf2, "  Supports both evolutionary and greedy search methods", 13, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: Training & Evaluation
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "Training & Evaluation", 32, ACCENT, True)

# Training
add_text_box(slide, 0.8, 1.3, 5.5, 0.5, "4 Training Pipelines", 20, ACCENT2, True)
trains = [
    ("Pretraining", "Standard dLLM pretraining\nELBO (MDLM), score-entropy (SEDD), VLB+CE (D3PM)"),
    ("DAG-Aware", "Biases masking by topological level\nCloses train-inference distribution gap"),
    ("Fine-tuning", "Answer-only loss\nFreeze backbone, train on answers"),
    ("Diffu-GRPO", "RL fine-tuning with GRPO\nOptimize for downstream task reward"),
]
for i, (name, desc) in enumerate(trains):
    y = 1.9 + i * 1.15
    add_rect(slide, 0.8, y, 2.0, 0.4, RGBColor(0x3A, 0x5A, 0x7A), name, 12)
    add_text_box(slide, 3.0, y - 0.05, 3.5, 1.0, desc, 11, LIGHT_GRAY)

# Evaluation
add_text_box(slide, 7.2, 1.3, 5.5, 0.5, "8 Benchmarks", 20, ACCENT2, True)
evals = [
    ("Code Gen", "MBPP, HumanEval\nMetric: pass@1"),
    ("Math", "GSM8K, MATH\nMetric: exact match"),
    ("QA", "HotpotQA\nMetric: exact match, F1"),
    ("Knowledge", "MMLU\nMetric: accuracy"),
    ("Logic", "ARC-Challenge, ProntoQA\nMetric: accuracy"),
]
for i, (name, desc) in enumerate(evals):
    y = 1.9 + i * 1.0
    add_rect(slide, 7.2, y, 1.8, 0.4, RGBColor(0x2D, 0x6E, 0x4A), name, 12)
    add_text_box(slide, 9.2, y - 0.05, 3.5, 0.9, desc, 11, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: Project Stats & Next Steps
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "Project Summary", 32, ACCENT, True)

# Stats table
stats = [
    ("Python Files", "69"),
    ("Lines of Code", "~9,200"),
    ("Config Files (YAML)", "31"),
    ("Test Cases", "60+"),
    ("Notebooks", "2"),
    ("CLI Entry Points", "5"),
    ("dLLM Models", "4"),
    ("Unmasking Schedulers", "5"),
    ("DAG Templates", "6"),
    ("Search Methods", "4"),
    ("Library Components", "9"),
    ("Ablation Configs", "7"),
]

add_text_box(slide, 0.8, 1.2, 4, 0.5, "By the Numbers", 20, ACCENT2, True)
for i, (k, v) in enumerate(stats):
    col = i // 6
    row = i % 6
    x = 0.8 + col * 3.2
    y = 1.8 + row * 0.5
    add_text_box(slide, x, y, 2.0, 0.4, k, 13, LIGHT_GRAY)
    add_text_box(slide, x + 2.0, y, 1.0, 0.4, v, 13, ACCENT3, True)

# CLI
add_text_box(slide, 7.5, 1.2, 5, 0.5, "CLI Commands", 20, ACCENT2, True)
cmds = [
    "dllm-train      # Train models",
    "dllm-eval       # Evaluate with schedulers",
    "dllm-eval-dags  # LLaDA + DAG strategies",
    "dllm-search     # Search optimal DAGs",
    "dllm-viz        # Visualize DAG structures",
]
tf = add_text_box(slide, 7.5, 1.8, 5, 3.0, "", 13)
for cmd in cmds:
    add_bullet(tf, cmd, 12, ACCENT3)

# Quick start
add_text_box(slide, 7.5, 4.2, 5, 0.5, "Quick Start", 20, ACCENT2, True)
tf2 = add_text_box(slide, 7.5, 4.8, 5, 2.0, "", 12)
add_bullet(tf2, "pip install -e '.[dev,library]'", 12, ACCENT3)
add_bullet(tf2, "dllm-train --model mdlm --dataset gsm8k", 12, LIGHT_GRAY)
add_bullet(tf2, "dllm-eval-dags --dags cot skeleton", 12, LIGHT_GRAY)
add_bullet(tf2, "dllm-search --method evolutionary", 12, LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
output_path = "docs/dLLM_Reason_V1.0.pptx"
prs.save(output_path)
print(f"Saved presentation to {output_path}")
