"""Generate v1.2.3 feature-summary presentation.

Covers:
  - v1.2.3 new features
  - LLaDA inference fixes
  - Sampling algorithm
  - All 8 unmasking strategies
  - Every parameter with default + CLI flag + how to set
  - Every script in scripts/ and scripts/runs/
  - Output saving feature

Usage:  python docs/generate_ppt_v1.2.3.py
Output: docs/dLLM_Reason_V1.2.3.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (same family as V1.0) ─────────────────────────────────────────────
BG       = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID   = RGBColor(0x22, 0x22, 0x3A)
CARD     = RGBColor(0x2A, 0x2A, 0x42)
ACCENT   = RGBColor(0x4E, 0xC9, 0xB0)   # teal
BLUE     = RGBColor(0x56, 0x9C, 0xD6)   # cornflower blue
YELLOW   = RGBColor(0xDC, 0xDC, 0xAA)  # warm yellow
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xCC, 0xCC, 0xCC)
MGRAY    = RGBColor(0x88, 0x88, 0x99)
ORANGE   = RGBColor(0xE0, 0x8A, 0x40)
RED      = RGBColor(0xC0, 0x50, 0x50)
GREEN    = RGBColor(0x50, 0xB0, 0x70)
PURPLE   = RGBColor(0x90, 0x60, 0xC0)

# ── Helpers ────────────────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def tb(slide, l, t, w, h, text, size=13, color=WHITE,
       bold=False, align=PP_ALIGN.LEFT, italic=False, name="Segoe UI"):
    """Add a text box. Returns text_frame."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = name
    p.alignment = align
    return tf


def para(tf, text, size=12, color=LGRAY, bold=False, space=4, italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = "Segoe UI"
    p.space_before = Pt(space)
    return p


def bullet(tf, text, size=12, color=LGRAY, level=0, space=3, mono=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Consolas" if mono else "Segoe UI"
    p.level = level
    p.space_before = Pt(space)
    return p


def rect(slide, l, t, w, h, fill, text="", size=12,
         fc=WHITE, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = fc
        p.font.bold = True
        p.font.name = "Segoe UI"
    return shape


def rr(slide, l, t, w, h, fill, text="", size=12, fc=WHITE):
    """Rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = fc
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    return shape


def hbar(slide, color=ACCENT, y=0, h=0.08):
    """Thin top accent bar."""
    rect(slide, 0, y, 13.333, h, color)


def footer(slide, text="dLLM-Reason v1.2.3"):
    rect(slide, 0, 7.3, 13.333, 0.2, BG_MID)
    tb(slide, 0.4, 7.3, 12, 0.2, text, 10, MGRAY)


def section_title(slide, text, color=ACCENT):
    tb(slide, 0.6, 0.15, 12, 0.65, text, 32, color, True)


# ── Table helper ───────────────────────────────────────────────────────────────

def table_row(slide, cols, y, row_h, widths, x_start, colors, sizes,
              fc_list, bold_list, mono_list=None):
    """Draw one row of a manual table."""
    x = x_start
    for i, (text, w) in enumerate(zip(cols, widths)):
        bg = colors[i] if isinstance(colors, list) else colors
        fc = fc_list[i] if isinstance(fc_list, list) else fc_list
        sz = sizes[i] if isinstance(sizes, list) else sizes
        bd = bold_list[i] if isinstance(bold_list, list) else bold_list
        mono = mono_list[i] if mono_list and isinstance(mono_list, list) else False

        rect(slide, x, y, w, row_h, bg)
        box = slide.shapes.add_textbox(
            Inches(x + 0.05), Inches(y + 0.03),
            Inches(w - 0.1), Inches(row_h - 0.06)
        )
        tf2 = box.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(sz)
        p.font.color.rgb = fc
        p.font.bold = bd
        p.font.name = "Consolas" if mono else "Segoe UI"
        p.alignment = PP_ALIGN.LEFT
        x += w


# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, ACCENT, 0, 0.12)

tb(s, 1.5, 1.3, 10, 1.3, "dLLM-Reason", 60, ACCENT, True, PP_ALIGN.CENTER)
tb(s, 1.5, 2.7, 10, 0.8,
   "DAG-Guided Discrete Diffusion LMs for Reasoning",
   26, WHITE, False, PP_ALIGN.CENTER)
rr(s, 5.4, 3.7, 2.5, 0.6, BLUE, "v1.2.3  Feature Summary", 20)

tf = tb(s, 1.5, 4.6, 10, 1.5, "", 14, LGRAY, False, PP_ALIGN.CENTER)
tf.paragraphs[0].text = "LLaDA Inference Fixes  ·  Official Sampling Algorithm"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = LGRAY
tf.paragraphs[0].font.name = "Segoe UI"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
para(tf, "Unified Eval Pipeline  ·  8 Unmasking Strategies", 14, LGRAY, False, 4)
para(tf, "Per-Sample Output Saving (JSON + Excel + Trajectory)", 14, LGRAY, False, 4)

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — v1.2.3 New Features Overview
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, ACCENT)
section_title(s, "v1.2.3  New Features")

cards = [
    (RED,    "🔧 LLaDA Inference Fixed",
     ["Correct mask_token_id from model.config",
      "apply_chat_template(return_tensors='pt') — no double-BOS",
      "System prompt prepended to user message (no 'system' role)",
      "temperature=0 → direct argmax, not logits/1e-6"]),
    (BLUE,   "⚡ Official Sampling Algorithm",
     ["Gumbel noise + argmax (matches GSAI-ML reference)",
      "Block-wise denoising: gen_length / block_length blocks",
      "Confidence-based remasking (low_confidence | random)",
      "Optional CFG support (cfg_scale > 0)"]),
    (GREEN,  "📊 Unified Eval Pipeline",
     ["YAML config as single source of truth",
      "CLI flags always override config",
      "8 benchmarks  ·  8 DAG strategies",
      "resume, verbose_errors, run_tests flags"]),
    (ORANGE, "💾 Per-Sample Output Saving",
     ["--save_outputs master switch",
      "QA pairs + ground truth per sample",
      "JSON + Excel (openpyxl) output files",
      "Optional trajectory recording (--record_trajectory)"]),
]

for i, (color, title, items) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.05 + row * 2.95
    rr(s, x, y, 6.2, 2.75, CARD)
    rect(s, x, y, 6.2, 0.52, color)
    tb(s, x + 0.15, y + 0.05, 5.9, 0.45, title, 16, WHITE, True)
    tf = tb(s, x + 0.15, y + 0.6, 5.9, 2.0, "", 13)
    tf.paragraphs[0].text = ""  # placeholder
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
            p.text = "▸  " + item
        else:
            p = tf.add_paragraph()
            p.text = "▸  " + item
        p.font.size = Pt(13)
        p.font.color.rgb = LGRAY
        p.font.name = "Segoe UI"
        p.space_before = Pt(4)

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — LLaDA Inference Fixes
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, RED)
section_title(s, "LLaDA Inference: Critical Fixes", RED)

tb(s, 0.6, 0.88, 12, 0.38,
   "All four bugs caused silent wrong output (wrong tokens, all-EOS, empty string). "
   "Fixed by matching the GSAI-ML reference implementation.",
   12, MGRAY)

fixes = [
    ("mask_token_id wrong",
     "tokenizer.mask_token_id → <|eot_id|> (id 128009)\nor <|startoftext|> (id 126080)",
     "Read from model.config.mask_token_id\n→ <|mdm_mask|> (id 126347)  ✓"),
    ("Double-BOS encoding",
     "tokenizer.encode(\n  apply_chat_template(tokenize=False)\n) → prepended BOS twice",
     "apply_chat_template(\n  messages, return_tensors='pt'\n)  ✓  — direct tensor"),
    ("System prompt malformed",
     "messages=[{'role':'system', ...}]\n→ LLaDA tokenizer ignores system role,\n  chat template broken",
     "user_content = f'{sys}\\n\\n{prompt}'\nmessages=[{'role':'user', 'content':user_content}]  ✓"),
    ("temperature=0 overflow",
     "logits / max(temperature, 1e-6)\n→ at temperature=0, division by 1e-6\n  → NaN/overflow in logits",
     "if temperature == 0:\n    return logits.argmax()  # direct\nelse: add Gumbel noise  ✓"),
]

for i, (title, bug_text, fix_text) in enumerate(fixes):
    y = 1.35 + i * 1.52
    # Title pill
    rr(s, 0.5, y, 12.3, 0.35, RGBColor(0x35, 0x20, 0x20), title, 12, RED)
    # Bug panel
    rect(s, 0.5, y + 0.38, 5.9, 1.05, RGBColor(0x28, 0x18, 0x18))
    tb(s, 0.55, y + 0.38, 1.1, 0.35, "BUG", 10, RED, True)
    tb(s, 0.6, y + 0.7, 5.7, 0.8, bug_text, 11, RGBColor(0xF0, 0x80, 0x80),
       name="Consolas")
    # Arrow
    tb(s, 6.5, y + 0.55, 0.5, 0.7, "→", 22, ACCENT, True, PP_ALIGN.CENTER)
    # Fix panel
    rect(s, 7.0, y + 0.38, 5.8, 1.05, RGBColor(0x18, 0x28, 0x1E))
    tb(s, 7.05, y + 0.38, 1.1, 0.35, "FIX", 10, GREEN, True)
    tb(s, 7.1, y + 0.7, 5.6, 0.8, fix_text, 11, RGBColor(0x80, 0xE0, 0xA0),
       name="Consolas")

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Official Sampling Algorithm
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, BLUE)
section_title(s, "Official LLaDA Block-Wise Sampling", BLUE)

# Left: algorithm steps
steps = [
    ("1", "Fill generation area with <|mdm_mask|> tokens",                   BLUE),
    ("2", "Divide into num_blocks = max_new_tokens / block_length blocks",    BLUE),
    ("3", "For each block, run steps_per_block sub-steps:",                   ACCENT),
    ("4", "  Forward pass → logits (shape: B × L × V)",                      ACCENT),
    ("5", "  Add Gumbel noise scaled by temperature → x̂₀ = argmax",          ACCENT),
    ("6", "  Compute confidence = softmax(logits)[x̂₀]",                     ACCENT),
    ("7", "  Commit top-k positions (highest confidence) in this block",      GREEN),
    ("8", "  Remask low-confidence tokens (low_confidence mode)",             ORANGE),
    ("9", "Force-fill any remaining masks at end (safety net)",               MGRAY),
]

for i, (num, text, color) in enumerate(steps):
    y = 1.0 + i * 0.68
    rr(s, 0.5, y + 0.08, 0.38, 0.38, color, num, 14)
    tb(s, 1.05, y, 7.2, 0.65, text, 13, LGRAY if num not in ("1","2","9") else WHITE,
       name="Consolas" if text.startswith("  ") else "Segoe UI")

# Right: block diagram
rect(s, 8.6, 0.95, 4.4, 6.15, CARD)
tb(s, 8.7, 1.0, 4.2, 0.45, "Visual: 3 blocks, 4 steps each", 13, ACCENT, True,
   PP_ALIGN.CENTER)

block_colors = [BLUE, ACCENT, GREEN]
block_names = ["Block 0  (steps 0–3)", "Block 1  (steps 4–7)", "Block 2  (steps 8–11)"]
step_labels = ["step 0", "step 1", "step 2", "step 3"]

for b, (bcolor, bname) in enumerate(zip(block_colors, block_names)):
    by = 1.6 + b * 1.9
    rect(s, 8.7, by, 4.2, 0.3, bcolor)
    tb(s, 8.75, by, 4.1, 0.3, bname, 11, WHITE, True)
    for st in range(4):
        sx = 8.75 + st * 1.0
        sy = by + 0.38
        filled = (b == 0) or (b == 1 and st < 3) or (b == 2 and st < 2)
        rect(s, sx, sy, 0.9, 1.1, bcolor if filled else RGBColor(0x30, 0x30, 0x48))
        tb(s, sx + 0.03, sy + 0.03, 0.84, 0.22, step_labels[st], 9, WHITE if filled else MGRAY)
        # Mini token cells
        for tk in range(6):
            tx = sx + 0.04 + tk * 0.138
            ty = sy + 0.32
            tok_filled = filled and (tk < (4 + st))
            rect(s, tx, ty, 0.12, 0.65,
                 bcolor if tok_filled else RGBColor(0x25, 0x25, 0x40))

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — 8 Unmasking Strategies
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, ACCENT)
section_title(s, "8 Unmasking Strategies")

tb(s, 0.6, 0.85, 12, 0.35,
   "All available via --dags flag or in configs/eval_default.yaml under dags:  dags:",
   12, MGRAY, False, PP_ALIGN.LEFT, True, "Consolas")

strategies = [
    ("confidence",    BLUE,   "scripts/runs/confidence.sh",
     "Highest-confidence tokens first at each step.\n"
     "LLaDA's default behaviour — strong baseline.\n"
     "No DAG object; uses ConfidenceScheduler."),
    ("random",        PURPLE, "scripts/runs/random.sh",
     "Uniform random selection from masked positions.\n"
     "Stochastic baseline for comparison.\n"
     "Uses RandomScheduler."),
    ("linear",        ACCENT, "scripts/runs/linear.sh",
     "Left-to-right sequential unmasking.\n"
     "Mimics autoregressive generation order.\n"
     "Uses LinearScheduler."),
    ("empty",         MGRAY,  "scripts/runs/empty.sh",
     "No structural constraint — purely random order.\n"
     "Standard LLaDA generation as-is.\n"
     "Uses RandomScheduler (same as 'random')."),
    ("cot",           ORANGE, "scripts/runs/cot.sh",
     "Chain-of-Thought DAG: cot_steps segments.\n"
     "Each segment depends on all prior segments.\n"
     "cot_steps default: 4  (--cot_steps N to change)."),
    ("skeleton",      GREEN,  "scripts/runs/skeleton.sh",
     "Skeleton-then-Detail: every 3rd token is 'structural'.\n"
     "Structural positions unmask before detail positions.\n"
     "Good for code / structured outputs."),
    ("bidirectional", BLUE,   "scripts/runs/bidirectional.sh",
     "4 segments; both ends unmask toward the center.\n"
     "Captures global context from both directions.\n"
     "Uses DAGScheduler + bidirectional_dag template."),
    ("answer_first",  RED,    "scripts/runs/answer_first.sh",
     "Last 20% of positions = answer region.\n"
     "Answer tokens unmask before reasoning chain.\n"
     "Useful for benchmarks where answer is a short span."),
]

for i, (name, color, script, desc) in enumerate(strategies):
    col = i % 4
    row = i // 4
    x = 0.35 + col * 3.22
    y = 1.28 + row * 2.95
    rr(s, x, y, 3.05, 2.78, CARD)
    rect(s, x, y, 3.05, 0.45, color)
    tb(s, x + 0.1, y + 0.05, 2.85, 0.38, name, 15, WHITE, True)
    tb(s, x + 0.1, y + 0.52, 2.85, 0.28, script, 9, YELLOW, name="Consolas")
    tb(s, x + 0.1, y + 0.85, 2.85, 1.78, desc, 11, LGRAY)

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Parameter Reference Part 1: Model + Inference + DAG
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, YELLOW)
section_title(s, "Parameter Reference  (1/2) — Model · Inference · DAG", YELLOW)
tb(s, 0.6, 0.88, 12, 0.32,
   "Set in configs/eval_default.yaml  ·  Override with CLI flag  ·  CLI always wins",
   11, MGRAY, False, PP_ALIGN.LEFT, True)

# Column widths: section | param | default | cli flag | choices / notes
COLW  = [1.3, 2.1, 2.2, 2.5, 4.7]
X0    = 0.35
RH    = 0.36
HBGS  = [RGBColor(0x2A, 0x2A, 0x50)] * 5
HFCS  = [ACCENT, ACCENT, YELLOW, BLUE, LGRAY]

# Header
table_row(s, ["Section", "Parameter", "Default", "CLI Flag", "Choices / Notes"],
          1.25, RH, COLW, X0,
          [RGBColor(0x18, 0x30, 0x50)] * 5,
          [11] * 5, HFCS, [True] * 5)

rows = [
    # section, param, default, cli, notes
    ("Model",     "model_id",      "checkpoints/llada-instruct", "--model_id",      "Local path or HF model ID"),
    ("Model",     "torch_dtype",   "bfloat16",                   "--torch_dtype",   "bfloat16 | float16 | float32"),
    ("Model",     "device_map",    "auto",                       "--device_map",    "auto | cuda:0 | cpu"),
    ("Inference", "num_steps",     "128",                        "--num_steps",     "Total diffusion denoising steps"),
    ("Inference", "block_length",  "32",                         "--block_length",  "Tokens per block; max_new_tokens ÷ block_length must be int"),
    ("Inference", "temperature",   "0.0",                        "--temperature",   "0 = greedy argmax; >0 = Gumbel sampling"),
    ("Inference", "cfg_scale",     "0.0",                        "--cfg_scale",     "Classifier-free guidance scale; 0 = disabled"),
    ("Inference", "remasking",     "low_confidence",             "--remasking",     "low_confidence | random"),
    ("Inference", "max_new_tokens","128",                        "--max_new_tokens","Total tokens to generate (alias: --generation_len)"),
    ("DAG",       "dags",          "[confidence]",               "--dags",          "Space-separated list; 8 choices (see slide 5)"),
    ("DAG",       "cot_steps",     "4",                          "--cot_steps",     "Reasoning segments for cot DAG"),
    ("DAG",       "mmlu_subjects", "null (10 defaults)",         "--mmlu_subjects", "Space-separated MMLU subject names"),
]

section_colors = {
    "Model":     RGBColor(0x20, 0x28, 0x40),
    "Inference": RGBColor(0x18, 0x28, 0x38),
    "DAG":       RGBColor(0x22, 0x20, 0x38),
}
alt = False
for (sec, param, default, cli, notes) in rows:
    y = 1.61 + rows.index((sec, param, default, cli, notes)) * RH
    bg = section_colors[sec]
    alt_bg = RGBColor(bg[0] + 8, bg[1] + 8, bg[2] + 8)
    row_bg = alt_bg if alt else bg
    alt = not alt
    table_row(s,
              [sec, param, default, cli, notes],
              y, RH, COLW, X0,
              [row_bg] * 5,
              [11, 11, 11, 11, 11],
              [BLUE, ACCENT, YELLOW, GREEN, LGRAY],
              [True, False, False, False, False],
              [False, True, True, True, False])

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Parameter Reference Part 2: Benchmarks + Output + Save
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, ORANGE)
section_title(s, "Parameter Reference  (2/2) — Benchmarks · Output · Save", ORANGE)
tb(s, 0.6, 0.88, 12, 0.32,
   "Flags that are 'store_true' (boolean) are listed as --flag to enable / --no_flag to disable",
   11, MGRAY, False, PP_ALIGN.LEFT, True)

table_row(s, ["Section", "Parameter", "Default", "CLI Flag", "Choices / Notes"],
          1.25, RH, COLW, X0,
          [RGBColor(0x18, 0x30, 0x50)] * 5,
          [11] * 5, HFCS, [True] * 5)

rows2 = [
    ("Benchmark", "benchmarks",        "[mbpp, humaneval]", "--benchmarks",          "mbpp humaneval hotpotqa mmlu gsm8k math arc prontoqa"),
    ("Benchmark", "num_samples",       "null (all)",        "--num_samples N",       "Integer; limits samples per benchmark"),
    ("Benchmark", "run_tests",         "true",              "--no_run_tests",        "Set flag to SKIP code execution (inspect output only)"),
    ("Benchmark", "verbose_errors",    "false",             "--verbose_errors",      "Print per-sample stderr / timeout / error message"),
    ("Output",    "output_dir",        "results",           "--output_dir PATH",     "Directory for result JSON files; .sh scripts append timestamp"),
    ("Output",    "resume",            "false",             "--resume",              "Skip runs whose result .json already exists"),
    ("Output",    "config",            "auto-detected",     "--config PATH",         "Load configs/eval_default.yaml; CLI overrides apply on top"),
    ("Save",      "save_outputs",      "false",             "--save_outputs",        "MASTER SWITCH — must be set for any files to be written"),
    ("Save",      "save_qa",           "true",              "--no_save_qa",          "Include prompt + generated answer in output files"),
    ("Save",      "save_ground_truth", "true",              "--no_save_ground_truth","Include reference/canonical answers in output files"),
    ("Save",      "record_trajectory", "false",             "--record_trajectory",   "Save per-step unmasking states (large file; use on small runs)"),
    ("Save",      "output_formats",    "[json, xlsx]",      "--output_formats ...",  "json | xlsx | json xlsx  (space-separated)"),
]

section_colors2 = {
    "Benchmark": RGBColor(0x20, 0x28, 0x20),
    "Output":    RGBColor(0x28, 0x20, 0x18),
    "Save":      RGBColor(0x28, 0x18, 0x28),
}
alt = False
for (sec, param, default, cli, notes) in rows2:
    y = 1.61 + rows2.index((sec, param, default, cli, notes)) * RH
    bg = section_colors2[sec]
    alt_bg = RGBColor(bg[0] + 8, bg[1] + 8, bg[2] + 8)
    row_bg = alt_bg if alt else bg
    alt = not alt
    table_row(s,
              [sec, param, default, cli, notes],
              y, RH, COLW, X0,
              [row_bg] * 5,
              [11, 11, 11, 11, 11],
              [ORANGE, ACCENT, YELLOW, GREEN, LGRAY],
              [True, False, False, False, False],
              [False, True, True, True, False])

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — scripts/ Reference
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, GREEN)
section_title(s, "scripts/  — Script Reference", GREEN)
tb(s, 0.6, 0.88, 12, 0.32,
   "All scripts load configs/eval_default.yaml by default. "
   "Extra CLI args override config values. Pass-through via \"$@\".",
   11, MGRAY, False, PP_ALIGN.LEFT, True)

scripts_top = [
    ("scripts/run_eval.sh",
     "Main LLaDA evaluation entry point.",
     ["Checks environment (PyTorch, dllm_reason package)",
      "Loads configs/eval_default.yaml automatically",
      "Appends timestamp to output_dir",
      "Passes all extra args to eval_dags.py",
      "Usage:  bash scripts/run_eval.sh  [--flag value ...]"],
     ACCENT),
    ("scripts/eval_dags.py",
     "Python evaluation driver — all logic lives here.",
     ["Parses CLI args with config-file fallback",
      "Loads LLaDA model, builds scheduler per DAG strategy",
      "Runs each benchmark × each strategy combination",
      "Saves per-run .json and summary.json",
      "Called by every .sh script; rarely invoked directly"],
     BLUE),
    ("scripts/infer_llada.py",
     "Standalone single-prompt inference — no eval harness.",
     ["Accepts --prompt directly on CLI",
      "Robust mask_token_id resolution (model.config → tokenizer)",
      "Prints decoded generation to stdout",
      "Usage:  python scripts/infer_llada.py --prompt \"...\""],
     PURPLE),
    ("scripts/download_models.py",
     "Download LLaDA weights from HuggingFace.",
     ["Saves to checkpoints/llada-instruct/ by default",
      "--mirror https://hf-mirror.com  for China",
      "Skips if directory already exists"],
     ORANGE),
    ("scripts/download_datasets.py",
     "Download benchmark datasets.",
     ["Saves to datasets/  by default",
      "--mirror flag for HF mirror",
      "Skips already-downloaded splits"],
     GREEN),
]

for i, (name, subtitle, items, color) in enumerate(scripts_top):
    col = i % 2 if i < 4 else 0
    row = i // 2
    if i == 4:
        x, y, w = 0.35, 6.1, 6.3
    else:
        x = 0.35 + col * 6.55
        y = 1.28 + row * 2.3
        w = 6.3
    h = 2.1 if i < 4 else 1.15

    rr(s, x, y, w, h, CARD)
    rect(s, x, y, w, 0.38, color)
    tb(s, x + 0.1, y + 0.04, w - 0.2, 0.32, name, 12, WHITE, True, name="Consolas")
    tb(s, x + 0.1, y + 0.44, w - 0.2, 0.3, subtitle, 11, YELLOW, True)
    tf = tb(s, x + 0.1, y + 0.78, w - 0.2, h - 0.85, "", 11)
    tf.paragraphs[0].text = ""
    for j, item in enumerate(items if i < 4 else items[:2]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "▸  " + item
        p.font.size = Pt(11)
        p.font.color.rgb = LGRAY
        p.font.name = "Segoe UI"
        p.space_before = Pt(2)

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — scripts/runs/ — Per-Strategy Scripts
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, BLUE)
section_title(s, "scripts/runs/  — Per-Strategy Scripts", BLUE)
tb(s, 0.6, 0.88, 12, 0.32,
   "Each script sets --dags <name> and --output_dir results/<name>_<timestamp>. "
   "All extra CLI args pass through: bash scripts/runs/cot.sh --num_samples 50",
   11, MGRAY, False, PP_ALIGN.LEFT, True)

run_scripts = [
    ("confidence.sh",     BLUE,   "Highest-confidence first (LLaDA default)"),
    ("random.sh",         PURPLE, "Uniform random unmasking (stochastic baseline)"),
    ("linear.sh",         ACCENT, "Left-to-right sequential unmasking"),
    ("empty.sh",          MGRAY,  "No constraint — pure standard LLaDA"),
    ("cot.sh",            ORANGE, "Chain-of-Thought DAG   (--cot_steps N)"),
    ("skeleton.sh",       GREEN,  "Skeleton-then-Detail DAG"),
    ("bidirectional.sh",  BLUE,   "Bidirectional DAG (both ends → center)"),
    ("answer_first.sh",   RED,    "Answer region unmasked before reasoning"),
    ("all_strategies.sh", YELLOW, "All 8 strategies in one run → single summary.json"),
    ("save_outputs.sh",   ACCENT, "Any strategy + --save_outputs + commented options"),
]

for i, (name, color, desc) in enumerate(run_scripts):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.28 + row * 1.18

    rr(s, x, y, 6.3, 1.05, CARD)
    rect(s, x, y, 0.25, 1.05, color)
    tb(s, x + 0.35, y + 0.08, 3.2, 0.4, name, 13, WHITE, True, name="Consolas")
    tb(s, x + 0.35, y + 0.52, 5.8, 0.48, desc, 12, LGRAY)

    # Usage snippet
    extra = " --cot_steps 6" if "cot" in name else ""
    usage = f"bash scripts/runs/{name}{extra}"
    tb(s, x + 3.65, y + 0.08, 2.55, 0.4, usage, 9, YELLOW, name="Consolas")

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Eval Pipeline & Config Flow
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, ACCENT)
section_title(s, "Eval Pipeline & Config Flow")

# Config priority diagram (top)
boxes_cfg = [
    (0.5,  2.2, 2.8, 1.1, "CLI Flags\n(highest priority)",                        ORANGE,  "1"),
    (4.0,  2.2, 2.8, 1.1, "configs/eval_default.yaml\n(loaded automatically)",    BLUE,    "2"),
    (7.5,  2.2, 2.8, 1.1, "argparse defaults\n(fallback)",                        MGRAY,   "3"),
    (10.7, 2.2, 2.0, 1.1, "eval_dags.py",                                        GREEN,   "→"),
]
for (x, y, w, h, text, color, badge) in boxes_cfg:
    rr(s, x, y, w, h, color)
    tb(s, x + 0.12, y + 0.1, w - 0.24, h - 0.15, text, 13, WHITE, True, PP_ALIGN.CENTER)

for ax in [3.3, 6.8, 10.2]:
    tb(s, ax, 2.55, 0.7, 0.42, "→", 22, ACCENT, True, PP_ALIGN.CENTER)

# YAML sections breakdown
tb(s, 0.5, 3.55, 12.4, 0.42,
   "configs/eval_default.yaml  sections:", 14, ACCENT, True)

yaml_sections = [
    ("model:",     "model_id  ·  torch_dtype  ·  device_map",      BLUE),
    ("inference:", "num_steps  ·  block_length  ·  temperature  ·  cfg_scale  ·  remasking  ·  max_new_tokens", ACCENT),
    ("benchmarks:","benchmarks  ·  num_samples  ·  run_tests  ·  verbose_errors", GREEN),
    ("dags:",      "dags  ·  cot_steps  ·  mmlu_subjects",          ORANGE),
    ("output:",    "output_dir  ·  resume",                         YELLOW),
    ("save:",      "save_outputs  ·  save_qa  ·  save_ground_truth  ·  record_trajectory  ·  output_formats", PURPLE),
]
for i, (section, keys, color) in enumerate(yaml_sections):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.28
    y = 4.05 + row * 0.88
    rr(s, x, y, 4.1, 0.78, CARD)
    rect(s, x, y, 4.1, 0.28, color)
    tb(s, x + 0.08, y + 0.02, 1.5, 0.26, section, 11, WHITE, True, name="Consolas")
    tb(s, x + 0.08, y + 0.33, 3.9, 0.4, keys, 10, LGRAY, name="Consolas")

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Output Saving Feature
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, ORANGE)
section_title(s, "Per-Sample Output Saving", ORANGE)

# Master switch callout
rr(s, 0.5, 0.9, 12.3, 0.62, RGBColor(0x2A, 0x20, 0x10))
tb(s, 0.65, 0.9, 12.0, 0.62,
   "Master switch:  --save_outputs   (or set  save_outputs: true  in configs/eval_default.yaml)",
   14, YELLOW, True, name="Consolas")

# 3 output file cards
file_cards = [
    ("{bench}_{dag}_samples.json", BLUE,
     "Full per-sample JSON records",
     ["prompt (if save_qa: true)",
      "generated output (raw + extracted)",
      "ground truth / canonical solution (if save_ground_truth: true)",
      "pass / fail status, error, stdout, stderr",
      "timed_out flag"]),
    ("{bench}_{dag}_samples.xlsx", GREEN,
     "Excel spreadsheet (openpyxl)",
     ["One row per sample",
      "Same columns as JSON, human-readable",
      "Long strings truncated to 2000 chars/cell",
      "Auto-fit column widths",
      "Easy to open in Excel / Google Sheets"]),
    ("{bench}_{dag}_trajectory.json", ORANGE,
     "Per-step unmasking states\n(--record_trajectory required)",
     ["Enabled with --record_trajectory",
      "One entry per sample",
      "List of decoded strings: one per diffusion step",
      "Shows generation area only (not prompt)",
      "⚠ Large file for long runs — keep num_samples small"]),
]

for i, (fname, color, subtitle, items) in enumerate(file_cards):
    x = 0.4 + i * 4.28
    rr(s, x, 1.7, 4.1, 5.45, CARD)
    rect(s, x, 1.7, 4.1, 0.42, color)
    tb(s, x + 0.1, 1.72, 3.9, 0.38, fname, 11, WHITE, True, name="Consolas")
    tb(s, x + 0.1, 2.18, 3.9, 0.35, subtitle, 11, YELLOW, True)
    tf = tb(s, x + 0.1, 2.6, 3.9, 4.2, "", 11)
    tf.paragraphs[0].text = ""
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "▸  " + item
        p.font.size = Pt(11)
        p.font.color.rgb = LGRAY
        p.font.name = "Segoe UI"
        p.space_before = Pt(5)

# Sub-options
tb(s, 0.5, 7.0, 12.4, 0.38,
   "Sub-options:  --no_save_qa   --no_save_ground_truth   --record_trajectory   --output_formats json   --output_formats xlsx   --output_formats json xlsx",
   10, MGRAY, name="Consolas")

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Quick Start Examples
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
hbar(s, ACCENT)
section_title(s, "Quick Start Examples")

examples = [
    ("Default run  (LLaDA + confidence scheduler)",
     "bash scripts/run_eval.sh",
     ACCENT),
    ("Override parameters on the fly",
     "bash scripts/run_eval.sh --benchmarks mbpp --num_samples 50 --num_steps 64",
     BLUE),
    ("Run a specific DAG strategy",
     "bash scripts/runs/cot.sh --benchmarks mbpp humaneval --cot_steps 6",
     ORANGE),
    ("Compare all 8 strategies",
     "bash scripts/runs/all_strategies.sh --num_samples 100",
     GREEN),
    ("Save QA pairs + ground truth to JSON + Excel",
     "bash scripts/run_eval.sh --save_outputs",
     YELLOW),
    ("Save outputs + record unmasking trajectory",
     "bash scripts/runs/save_outputs.sh --record_trajectory --num_samples 10",
     ORANGE),
    ("Run only JSON output (skip Excel)",
     "bash scripts/run_eval.sh --save_outputs --output_formats json",
     BLUE),
    ("Single-prompt inference (no eval harness)",
     'python scripts/infer_llada.py --model_id checkpoints/llada-instruct \\\n'
     '    --prompt "Write a Python function to reverse a string." \\\n'
     '    --num_steps 128 --temperature 0.0',
     PURPLE),
    ("Edit config file to change all defaults",
     "# edit configs/eval_default.yaml, then:\nbash scripts/run_eval.sh",
     MGRAY),
]

for i, (title, cmd, color) in enumerate(examples):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.3
    y = 1.05 + row * 2.1
    h = 1.9
    rr(s, x, y, 4.1, h, CARD)
    rect(s, x, y, 4.1, 0.3, color)
    tb(s, x + 0.1, y + 0.03, 3.9, 0.25, title, 10, WHITE, True)
    rect(s, x + 0.1, y + 0.36, 3.9, h - 0.46, RGBColor(0x12, 0x12, 0x22))
    tb(s, x + 0.18, y + 0.42, 3.72, h - 0.56, cmd, 10, YELLOW, name="Consolas")

footer(s)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out = "docs/dLLM_Reason_V1.2.3.pptx"
prs.save(out)
print(f"Saved: {out}")
